"""Render a deck outline (JSON on stdin) to a .pptx file via python-pptx.

Sidecar protocol (must match the other scripts in this directory):
- UTF-8 stdout regardless of the Windows console codepage;
- the verdict is the LAST non-empty stdout line, one JSON object:
    {"ok": true} or {"ok": false, "error": "..."}
- library warnings may print above the verdict; the Rust caller reads only
  the last non-empty line.

Input (stdin, UTF-8 JSON):
    {"path": "...", "title": "...", "subtitle": "...",
     "slides": [{"title": "...", "bullets": ["..."], "note": "..."}]}
"""

import json
import sys


def build(payload: dict) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    # Cover slide (default template's title layout).
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = payload.get("title", "")
    if payload.get("subtitle", ""):
        cover.placeholders[1].text = payload["subtitle"]

    for item in payload.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = item.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = item.get("bullets", [])
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(18)
            paragraph.level = 0
        note = item.get("note", "")
        if note:
            slide.notes_slide.notes_text_frame.text = note

    prs.save(payload["path"])


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    try:
        raw = sys.stdin.buffer.read().decode("utf-8")
        payload = json.loads(raw)
        build(payload)
        print(json.dumps({"ok": True}, ensure_ascii=False))
    except Exception as exc:  # noqa: BLE001 - verdict protocol needs one line
        print(json.dumps({"ok": False, "error": str(exc)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
